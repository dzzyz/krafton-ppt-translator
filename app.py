import streamlit as st
import re, json, time, io, os
from pptx import Presentation
from pptx.util import Pt
import anthropic
from datetime import datetime

# ── 페이지 설정 ────────────────────────────────────────────
st.set_page_config(
    page_title="KRAFTON PPT Translator",
    page_icon="🏢",
    layout="wide"
)

# ── 상수 ──────────────────────────────────────────────────
MODEL         = "claude-sonnet-4-20250514"
LANG_MAP      = {"English": "English", "Japanese": "Japanese", "Chinese": "Traditional Chinese"}
SKIP_NAMES    = ("slide number", "date", "footer", "page")
KOREAN_FONTS  = {
    "Pretendard", "나눔고딕", "맑은 고딕", "Malgun Gothic",
    "NanumGothic", "KoPubWorldBatang", "굴림", "돋움", "바탕", "나눔명조"
}
EN_FONT        = "Calibri"

# ── 기본 Glossary ──────────────────────────────────────────
BASE_GLOSSARY = {
    # 인명
    "장병규": "BG Chang",      "김창한": "CH Kim",         "배동근": "DK Bae",
    "장태석": "TS Jang",       "오진호": "Jin Oh",          "이강욱": "Kangwook Lee",
    "박혜리": "Maria Park",    "윤상훈": "Richard Yoon",   "손현일": "Sean Sohn",
    "한소영": "Soyoung Han",   "이병욱": "Andy Lee",        "김정연": "Jung Yun Kim",
    "김낙형": "Harns Kim",     "박찬민": "Chanmin Park",    "홍정택": "Jeongtack Hong",
    "박재민": "Jaemin Park",   "강재형": "Jaehyung Kang",  "박상훈": "Albert Park",
    "정용": "Young Chung",     "윤수진": "Sujin Yun",       "최승환": "Seunghwan Choi",
    "이동훈": "Donghun Lee",   "조혁일": "Hyuk Cho",        "유재영": "Jaeyoung Yoo",
    "황동경": "Dongkyung Hwang", "김현아": "Hyuna Kim",     "이석현": "Matthew Lee",
    "이창준": "Changjun Lee",  "이지은": "Jieun Lee",       "현종민": "Jongmin Hyun",
    "김정화": "Junghwa Kim",   "신소희": "Sohee Shin",      "노경원": "Kyoungwon Noh",
    "정지현": "Jeehyun Jung",  "김지영": "Jiyoung Kim",     "박서훈": "Seohoon Park",
    "김인영": "Inyoung Kim",   "김고운": "Kowoon Kim",      "어창선": "Chang Seon Eo",
    # 재무
    "매출": "Sales",
    "영업이익": "Operating Profit (OP)",
    "영업이익률": "OPM",
    "회계 매출": "Reported Revenue",
    "재무계획": "Financial Plan",
    "예산": "Budget",
    "연초 실적 점검": "Review of Early-year Performance",
    "전년 대비": "YoY",
    "누적 실적": "Cumulative Performance",
    "계획 대비": "vs. Plan",
    "인건비": "Labor Costs",
    "마케팅비": "Marketing Expenses",
    "지급수수료": "Service Fees",
    "임차료/상각비": "Rental/Depreciation",
    "비통제 비용": "Uncontrollable Costs",
    "매출연동비": "App Fees / Cost of Sales",
    "예비비": "Reserve Fund",
    "손익": "P/L",
    "범 크래프톤": "KRAFTON Family",
    "크래프톤 Business": "KRAFTON Business",
    "비게임 프로젝트": "Non-gaming Projects",
    "전사 비용": "Company-wide Expenses",
    "예산 집중 관리": "Intensive Budget Control",
    # 전략
    "핵심 서비스": "Core Service",
    "신작": "New IP",
    "신규 IP": "New IP",
    "장기 PLC화": "long-term PLC",
    "직접 서비스": "Direct Service",
    "두 자릿수 성장": "Double-digit Growth",
    "계단식 성장": "Stepwise Growth",
    "파이프라인": "Pipeline",
    "보고": "Report",
    "승인/보고": "Approval/Report",
    "승인사항": "Approval Item",
    "보류": "Hold (Deferred)",
    "서면 보고": "Written Report",
    "사전공유": "Pre-sharing",
    "서비스 종료": "Sunset / Discontinue",
    "퍼블리싱 운영 최적화 방안": "Publishing Operations Optimization Plan",
    "운영 Agility": "Operational Agility",
    "체질 개선": "Operational Transformation",
    "런치패드 프로그램": "Launchpad Program",
    "Tentpole 캠페인": "Tentpole Campaign",
    # 조직/HR
    "이사회": "Board of Directors (BOD)",
    "대표이사": "CEO",
    "이사회 의장": "Board Chair",
    "사외이사": "Outside Director",
    "제작 리더십": "Production Leadership",
    "제작 리더": "Production Lead",
    "제작총괄": "Head of Production",
    "비제작 조직": "Non-development Org",
    "신사업 자회사": "New Business Subsidiaries",
    "전문계약직": "Professional Contractors",
    "사내채용": "Internal Hiring",
    "인력계획": "Workforce Plan",
    "증원": "Headcount Increase",
    "유휴인력": "Idle Employees",
    "低성과자": "Low-performing Employees",
    "자발적 퇴직": "Voluntary Exit",
    "희망퇴직": "Voluntary Retirement",
    "정규직": "Full-time Employee",
    "미등기임원": "Non-executive Officers",
    "책임파트너제도": "Accountability Partner System",
    "별첨": "Separately Attached",
}

# ── Glossary DB 파일 핸들러 ────────────────────────────────
GLOSSARY_DB_FILE = "glossary_db.json"

def load_glossary_db():
    if not os.path.exists(GLOSSARY_DB_FILE):
        return {"approved_glossary": {}, "pending_glossary": []}
    try:
        with open(GLOSSARY_DB_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {"approved_glossary": {}, "pending_glossary": []}

def save_glossary_db(db):
    try:
        with open(GLOSSARY_DB_FILE, "w", encoding="utf-8") as f:
            json.dump(db, f, indent=2, ensure_ascii=False)
    except Exception as e:
        st.error(f"DB 저장 오류: {e}")

# ── Session state 초기화 (로컬변수 연동) ──────────────
if "session_extra_glossary" not in st.session_state:
    st.session_state.session_extra_glossary = {}
if "admin_logged_in"        not in st.session_state:
    st.session_state.admin_logged_in = False

# Translation Review session state
for _k, _v in {
    "tr_slides": [],          # [{ko, en, max_chars, en_len}]
    "tr_has_data": False,
    "tr_file_name": "",
}.items():
    if _k not in st.session_state:
        st.session_state[_k] = _v


def get_active_glossary():
    """BASE + 승인된 항목(DB) + 이번 세션 추가 = 최종 glossary"""
    db = load_glossary_db()
    merged = dict(BASE_GLOSSARY)
    merged.update(db.get("approved_glossary", {}))
    merged.update(st.session_state.session_extra_glossary)
    return merged


# ══════════════════════════════════════════════════════════
# PPT 처리 함수
# ══════════════════════════════════════════════════════════

def should_skip(shape):
    return any(k in getattr(shape, "name", "").lower() for k in SKIP_NAMES)


def has_korean(text):
    return bool(re.search(r"[\uAC00-\uD7A3]", text))


def iter_paragraphs(shapes):
    """shape_id + para_idx 기반 순회 — 그룹/표/중첩 그룹 모두 지원"""
    for shape in shapes:
        if getattr(shape, "has_text_frame", False) and shape.text_frame:
            for pi, para in enumerate(shape.text_frame.paragraphs):
                yield shape, para, pi
        if getattr(shape, "has_table", False) and shape.table:
            pi_counter = 0
            for row in shape.table.rows:
                for cell in row.cells:
                    if getattr(cell, "text_frame", None):
                        for para in cell.text_frame.paragraphs:
                            yield shape, para, pi_counter
                            pi_counter += 1
        if getattr(shape, "shape_type", None) == 6:
            yield from iter_paragraphs(shape.shapes)


def get_shape_width_pt(shape):
    try:
        w_emu = shape.width
        if w_emu:
            tf       = getattr(shape, "text_frame", None)
            left_in  = tf.margin_left  if (tf and tf.margin_left)  else 91440
            right_in = tf.margin_right if (tf and tf.margin_right) else 91440
            return max((w_emu - left_in - right_in) / 12700, 10)
    except Exception:
        pass
    return None


def get_slide_texts(slide):
    result = []
    for global_idx, (shape, para, pi) in enumerate(iter_paragraphs(slide.shapes)):
        if should_skip(shape):
            continue
        full = para.text
        if not full.strip() or not has_korean(full):
            continue
        font_pt = None
        for run in para.runs:
            if run.font.size:
                font_pt = run.font.size.pt
                break
        result.append({
            "shape_id":     shape.shape_id,
            "para_idx":     pi,
            "text":         full,
            "font_pt":      font_pt,
            "box_width_pt": get_shape_width_pt(shape),
        })
    return result


def detect_slide_type(texts):
    combined = " ".join(t["text"] for t in texts)
    if any(k in combined for k in ["승인사항", "승인 요청", "의결"]):
        return "approval"
    if any(k in combined for k in ["억원", "매출", "손익", "YoY"]):
        return "financial"
    if any(k in combined for k in ["출시", "파이프라인", "Q1", "Q2"]):
        return "timeline"
    if "사전공유" in combined:
        return "pre-disclosure"
    return "strategy"


def translate_slide(client, texts, slide_type, target_lang_str, glossary):
    if not texts:
        return {}

    type_rules = {
        "financial":      "Keep all numbers/units exact. Translate labels/headers only.",
        "approval":       'Use formal request language: "We request approval for..."',
        "strategy":       "Use concise noun phrases for bullets. Lead with the conclusion.",
        "timeline":       "Translate: 출시예정→Scheduled, 지연→Delayed, 완료→Released",
        "pre-disclosure": 'Add "(Pre-Disclosure)" prefix if not present.',
    }

    # Build input — no space constraints in translation (meaning first)
    input_map = {
        str(i): {"text": t["text"]}
        for i, t in enumerate(texts)
    }

    # Filter glossary to only terms that appear in this slide
    slide_text = " ".join(t["text"] for t in texts)
    relevant_glossary = {k: v for k, v in glossary.items() if k in slide_text}
    gstr = "\n".join(f"  {k} → {v}" for k, v in relevant_glossary.items()) if relevant_glossary else "(해당 슬라이드에 glossary 용어 없음)"

    prompt = f"""Translate the following Korean PowerPoint texts into professional {target_lang_str} for a KRAFTON Board of Directors meeting.

Slide type: {slide_type}
{type_rules.get(slide_type, "")}

## Mandatory glossary (use EXACTLY as shown):
{gstr}

## Do NOT translate or modify:
- Numbers, %, financial figures
- Game titles: PUBG, BGMI, OVERDARE, Black Budget, Valor, inZOI
- Company names: KRAFTON, Unknown Worlds, Neon Giant
- Labels: DRI, SL, □, ■, N/A, As-Is, To-Be

## Translation guidelines:
1. MEANING FIRST — preserve the original meaning, nuance, and tone COMPLETELY. This is a formal board meeting document. Do NOT omit or summarize any content.
2. Use natural, executive-level English that a native speaker would write for a board deck.
3. Do NOT over-shorten. If the original has detailed reasoning, keep it. Cutting meaning is worse than slightly longer text.
4. Preserve \\n line breaks exactly as in the original.
5. Preserve the exact position of (사전공유)→(Pre-sharing) within the sentence. Do not move it.
6. For financial slides: keep all numbers and units exact, translate labels with precision.
7. For approval slides: use formal language e.g. "We hereby request approval for..."

## Input (JSON — index → {{"text"}}):
{json.dumps(input_map, ensure_ascii=False, indent=2)}

## Output format:
Return ONLY a valid JSON object. Same index keys, string values only. No markdown, no explanation.
Example: {{"0": "Translated text", "1": "Another translation"}}"""

    system_prompt = "You are a professional Korean-English translator specializing in corporate board meeting materials for KRAFTON. Your top priority is accurate, natural, high-quality English that preserves the original meaning completely. Return ONLY valid JSON with string values."

    res = client.messages.create(
        model=MODEL,
        max_tokens=4096,
        system=system_prompt,
        messages=[{"role": "user", "content": prompt}]
    )
    raw   = res.content[0].text.strip()
    match = re.search(r"\{[\s\S]*\}", raw)
    if not match:
        raise ValueError("JSON 파싱 실패")
    matched_str = match.group()
    try:
        return json.loads(matched_str)
    except Exception:
        import ast
        try:
            safe_str = matched_str.replace("null", "None").replace("true", "True").replace("false", "False")
            return ast.literal_eval(safe_str)
        except Exception:
            raise ValueError("JSON 및 보조 파싱 모두 실패")


def replace_para_text(para, new_text, shape=None, min_pt=7):
    # ── 방어: string이 아니면 무조건 skip ─────────────────
    if not new_text or not isinstance(new_text, str):
        return
    new_text = new_text.strip()
    if not new_text:
        return

    # 첫 run에서 서식 추출
    orig_font_size = None
    orig_font_name = None
    for run in para.runs:
        if run.font.size:
            orig_font_size = run.font.size.pt
        if run.font.name:
            orig_font_name = run.font.name
        break

    # run 교체
    for i, run in enumerate(para.runs):
        if i == 0:
            run.text = new_text
            
            if orig_font_size:
                final_pt = orig_font_size
                min_allowed = max(orig_font_size - 4, min_pt)
                
                if shape is not None:
                    try:
                        w_emu = shape.width
                        if w_emu:
                            tf = getattr(shape, "text_frame", None)
                            left_in  = tf.margin_left  if (tf and getattr(tf, "margin_left", None))  else 91440
                            right_in = tf.margin_right if (tf and getattr(tf, "margin_right", None)) else 91440
                            box_w = max((w_emu - left_in - right_in) / 12700, 10)
                            lines = new_text.split('\n')
                            max_line_len = max(len(l) for l in lines) if lines else len(new_text)
                            if max_line_len > 0:
                                required_font = box_w / (max_line_len * 0.55)
                                if required_font < final_pt:
                                    final_pt = max(round(required_font, 1), min_allowed)
                    except Exception:
                        pass
                
                run.font.size = Pt(final_pt)

            if orig_font_name and orig_font_name in KOREAN_FONTS:
                run.font.name = EN_FONT
        else:
            run.text = ""

    # run이 없는 경우 새로 생성
    if not para.runs:
        from pptx.oxml.ns import qn
        from lxml import etree
        r   = etree.SubElement(para._p, qn("a:r"))
        etree.SubElement(r, qn("a:rPr"), attrib={"lang": "en-US"})
        t   = etree.SubElement(r, qn("a:t"))
        t.text = new_text


# ══════════════════════════════════════════════════════════
# UI
# ══════════════════════════════════════════════════════════
c_top1, c_top2 = st.columns([3, 1])
with c_top2:
    st.link_button("💬 Glossary 제안 / 의견 남기기", "https://docs.google.com/forms/d/e/1FAIpQLSezU-H6m0TMt2Ve-QUTZv483JklIdtfAsKi7rvYNW74l5B5lw/viewform", use_container_width=True)

tab_translate, tab_delta, tab_glossary = st.tabs([
    "🚀 번역", "🔄 Delta 번역", "📖 Glossary"
])


# ──────────────────────────────────────────────────────────
# TAB 1: 번역
# ──────────────────────────────────────────────────────────
with tab_translate:
    st.header("🏢 KRAFTON BOD PPT Translator")
    st.caption(
        "BOD 자료 PPT를 올려주시면 AI가 번역해드립니다 🙌 "
        "인명·용어 Glossary 자동 적용, 원문 의미 100% 보존 번역! "
        "텍스트 박스 크기와 위치는 원본 그대로 유지하고, 영문이 길어지면 폰트 크기만 자동 조정합니다. "
        "번역 후에는 슬라이드별 원문/번역 비교 + 📖 Glossary 미적용을 자동 감지해드려요."
    )

    col1, col2 = st.columns(2)
    with col1:
        target_lang = st.selectbox("번역 언어", ["English", "Japanese", "Chinese"])
    with col2:
        user_min_pt = st.number_input("최소 허용 폰트 크기 (pt)", min_value=1, max_value=40, value=7)


    # API Key — Streamlit Secrets에서 자동 로드
    api_key = st.secrets.get("ANTHROPIC_API_KEY", "")
    if not api_key:
        st.error("⚠️ API Key 미설정. Streamlit Cloud → Settings → Secrets에 ANTHROPIC_API_KEY를 추가해주세요.")

    active_glossary = get_active_glossary()
    st.caption(
        f"현재 적용 Glossary: **{len(active_glossary)}개** 항목 "
        f"(기본 {len(BASE_GLOSSARY)} "
        f"+ 승인 {len(load_glossary_db()['approved_glossary'])} "
        f"+ 세션 {len(st.session_state.session_extra_glossary)})"
    )

    uploaded_file = st.file_uploader("PPT 파일 업로드 (.pptx)", type=["pptx"])

    if uploaded_file and api_key:
        st.success(f"✅ **{uploaded_file.name}** 업로드 완료")

        if st.button("🚀 번역 시작", type="primary", use_container_width=True):
            start_time = time.time()
            client          = anthropic.Anthropic(api_key=api_key)
            target_lang_str = LANG_MAP.get(target_lang, "English")
            active_glossary = get_active_glossary()

            # 파싱
            pptx_bytes = uploaded_file.read()
            prs        = Presentation(io.BytesIO(pptx_bytes))

            with st.spinner("📊 텍스트 파싱 중..."):
                all_slides_info = []
                total_ko        = 0
                for slide in prs.slides:
                    texts = get_slide_texts(slide)
                    all_slides_info.append(texts)
                    total_ko += len(texts)

            st.info(f"총 {len(prs.slides)}장 · 번역 대상 **{total_ko}개** 단락")

            # 번역
            progress_bar     = st.progress(0)
            status_text      = st.empty()
            log_area         = st.empty()
            all_translations = []
            log_lines        = []

            for si, texts in enumerate(all_slides_info):
                ko_count = len(texts)
                progress_bar.progress((si + 1) / len(all_slides_info))

                if ko_count == 0:
                    all_translations.append({})
                    log_lines.append(f"⏭️  Slide {si+1:2d} — 한국어 없음")
                else:
                    slide_type = detect_slide_type(texts)
                    status_text.text(f"Slide {si+1}/{len(all_slides_info)} [{slide_type}] 번역 중...")
                    for attempt in range(2):
                        try:
                            translated = translate_slide(
                                client, texts, slide_type, target_lang_str, active_glossary
                            )
                            all_translations.append(translated)
                            log_lines.append(f"✅ Slide {si+1:2d} [{slide_type:15s}] — {ko_count}개")
                            break
                        except Exception as e:
                            if attempt == 0:
                                log_lines.append(f"⚠️  Slide {si+1:2d} 재시도... ({e})")
                                time.sleep(2)
                            else:
                                all_translations.append({})
                                log_lines.append(f"❌  Slide {si+1:2d} 실패 — 원본 유지")
                    time.sleep(0.3)

                log_area.code("\n".join(log_lines))

            # PPT 텍스트 교체 (shape_id + para_idx 기반)
            status_text.text("💾 PPT 생성 중...")
            for si, (slide, texts, translated_map) in enumerate(
                zip(prs.slides, all_slides_info, all_translations)
            ):
                if not translated_map:
                    continue

                # shape_id → shape 매핑 테이블 생성
                shape_map = {}
                def _collect_shapes(shapes):
                    for shape in shapes:
                        shape_map[shape.shape_id] = shape
                        if getattr(shape, "shape_type", None) == 6:
                            _collect_shapes(shape.shapes)
                _collect_shapes(slide.shapes)

                for ti, text_info in enumerate(texts):
                    tr = translated_map.get(str(ti))
                    if isinstance(tr, dict):
                        tr = tr.get("text") or tr.get("translation") or ""
                    
                    if not tr or not isinstance(tr, str):
                        continue
                    tr = tr.strip()
                    if not tr:
                        continue

                    shape_id = text_info["shape_id"]
                    para_idx = text_info["para_idx"]
                    shape    = shape_map.get(shape_id)
                    if shape is None:
                        continue

                    # 해당 para 찾기 (텍스트프레임 / 표 셀)
                    para = None
                    if getattr(shape, "has_text_frame", False) and shape.text_frame:
                        paras = shape.text_frame.paragraphs
                        if para_idx < len(paras):
                            para = paras[para_idx]
                    elif getattr(shape, "has_table", False) and shape.table:
                        all_paras = []
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if getattr(cell, "text_frame", None):
                                    all_paras.extend(cell.text_frame.paragraphs)
                        if para_idx < len(all_paras):
                            para = all_paras[para_idx]

                    if para is not None:
                        replace_para_text(para, tr, shape=shape, min_pt=user_min_pt)

            # 저장 & 다운로드
            output = io.BytesIO()
            prs.save(output)
            output.seek(0)
            out_name = uploaded_file.name.rsplit(".", 1)[0] + f"_{target_lang[:2].upper()}.pptx"

            progress_bar.progress(1.0)
            status_text.empty()
            elapsed_time = time.time() - start_time
            st.success(f"🎉 번역 완료! (⏱️ 소요시간: {elapsed_time:.1f}초)")
            st.download_button(
                label=f"⬇️ {out_name} 다운로드",
                data=output,
                file_name=out_name,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
                type="primary",
            )

            # ── 번역 결과를 session state에 저장 (검토용) ──
            tr_data = []
            for si, (texts, translated_map) in enumerate(zip(all_slides_info, all_translations)):
                pairs = []
                for ti, text_info in enumerate(texts):
                    orig = text_info["text"]
                    tr = translated_map.get(str(ti), "")
                    if isinstance(tr, dict):
                        tr = tr.get("text") or tr.get("translation") or ""
                    if not isinstance(tr, str):
                        tr = ""
                    # Calculate max_chars for display
                    font_pt = text_info.get("font_pt")
                    box_w = text_info.get("box_width_pt")
                    max_chars = None
                    if font_pt and box_w and font_pt > 0:
                        cpl = int(box_w / (font_pt * 0.52))
                        nl = orig.count("\n") + 1
                        max_chars = int(cpl * nl * 1.1)
                    pairs.append({
                        "ko": orig,
                        "en": tr.strip(),
                        "max_chars": max_chars,
                        "en_len": len(tr.strip()),
                    })
                tr_data.append(pairs)
            st.session_state.tr_slides = tr_data
            st.session_state.tr_has_data = True
            st.session_state.tr_file_name = uploaded_file.name

    # ══════════════════════════════════════════════════
    # 번역 결과 검토 (rule-based)
    # ══════════════════════════════════════════════════
    if st.session_state.tr_has_data:
        st.divider()
        st.subheader("📋 번역 결과 검토")
        st.caption(
            "슬라이드별로 원문과 번역을 나란히 확인할 수 있습니다. "
            "📖 **Glossary 미적용** — 등록된 용어가 번역에 반영되지 않은 경우를 자동 감지합니다."
        )

        tr_data = st.session_state.tr_slides
        glossary = get_active_glossary()

        # ── Rule-based checks (Glossary only) ──
        def check_slide(pairs, glossary):
            """Check glossary application. Case-insensitive, with partial-match guard."""
            warnings = []
            all_ko = " ".join(p["ko"] for p in pairs)
            all_en_lower = " ".join(p["en"] for p in pairs).lower()

            # Sort glossary by term length descending — longer terms take priority
            sorted_terms = sorted(glossary.items(), key=lambda x: len(x[0]), reverse=True)
            checked_ko_spans = []

            for ko_term, en_term in sorted_terms:
                if len(ko_term) < 2:
                    continue
                if ko_term not in all_ko:
                    continue
                # Skip if this term is a substring of an already-checked longer term
                is_substring = False
                for checked in checked_ko_spans:
                    if ko_term in checked and ko_term != checked:
                        is_substring = True
                        break
                if is_substring:
                    continue
                checked_ko_spans.append(ko_term)

                # Case-insensitive check
                if en_term.lower() not in all_en_lower:
                    warnings.append(("📖", f"Glossary: '{ko_term}' → '{en_term}' 미적용 가능성"))

            return warnings

        # Run checks
        all_checks = {}
        total_issues = 0
        for si, pairs in enumerate(tr_data):
            if pairs:
                checks = check_slide(pairs, glossary)
                all_checks[si] = checks
                total_issues += len(checks)

        # Summary
        slides_with_issues = sum(1 for c in all_checks.values() if c)
        sm1, sm2, sm3 = st.columns(3)
        with sm1:
            st.metric("전체 슬라이드", f"{len(tr_data)}장")
        with sm2:
            st.metric("⚠️ 확인 필요", f"{slides_with_issues}장")
        with sm3:
            st.metric("감지된 항목", f"{total_issues}건")

        # TXT export
        if tr_data:
            lines = ["═"*50, "  번역 결과 검토 Report", f"  {datetime.now().strftime('%Y-%m-%d %H:%M')}", "═"*50, ""]
            lines.append(f"  전체: {len(tr_data)}장 / 확인 필요: {slides_with_issues}장 / 감지: {total_issues}건")
            lines.append("")
            for si, pairs in enumerate(tr_data):
                if not pairs:
                    continue
                checks = all_checks.get(si, [])
                icon = "⚠️" if checks else "✅"
                lines.append("─"*50)
                lines.append(f"  슬라이드 {si+1}  {icon}")
                lines.append("─"*50)
                for p in pairs:
                    lines.append(f"  KR: {p['ko']}")
                    lines.append(f"  EN: {p['en']}")
                    lines.append("")
                if checks:
                    for tag, msg in checks:
                        lines.append(f"  {tag} {msg}")
                else:
                    lines.append("  ✅ 이상 없음")
                lines.append("")
            lines += ["═"*50, "  END OF REPORT", "═"*50]
            st.download_button("📝 검토 리포트 (TXT)", "\n".join(lines),
                               file_name=f"Translation_Review_{datetime.now().strftime('%Y%m%d')}.txt",
                               mime="text/plain", use_container_width=True, key="tr_dl_txt")

        # ── Per-slide display ──
        for si, pairs in enumerate(tr_data):
            if not pairs:
                continue
            checks = all_checks.get(si, [])

            # Header
            if checks:
                st.markdown(
                    f'<div style="display:flex;align-items:center;gap:10px;margin-top:16px;margin-bottom:8px;">'
                    f'<span style="font-size:15px;font-weight:700;color:#111827;">슬라이드 {si+1}</span>'
                    f'<span style="font-size:11px;font-weight:600;padding:3px 12px;border-radius:16px;'
                    f'background:#FFFBEB;color:#92400E;border:1px solid #FDE68A;">⚠️ {len(checks)}건 확인</span></div>',
                    unsafe_allow_html=True)
            else:
                st.markdown(
                    f'<div style="display:flex;align-items:center;gap:10px;margin-top:16px;margin-bottom:8px;">'
                    f'<span style="font-size:15px;font-weight:700;color:#111827;">슬라이드 {si+1}</span>'
                    f'<span style="font-size:11px;font-weight:600;padding:3px 12px;border-radius:16px;'
                    f'background:#F0FDF4;color:#166534;border:1px solid #BBF7D0;">✅ OK</span></div>',
                    unsafe_allow_html=True)

            # Clean text comparison table — no space column
            table_data = [{"원문 (한국어)": p["ko"], "번역 (English)": p["en"]} for p in pairs]
            st.dataframe(table_data, use_container_width=True, hide_index=True)

            # Warnings
            for tag, msg in checks:
                styles = {
                    "📖": "background:#EFF6FF;border-left:3px solid #3B82F6;color:#1E40AF",
                }
                st.markdown(
                    f'<div style="padding:6px 12px;margin:3px 0;border-radius:6px;font-size:12px;'
                    f'line-height:1.5;{styles.get(tag, "")}">{tag} {msg}</div>',
                    unsafe_allow_html=True)


# ──────────────────────────────────────────────────────────
# TAB 2: Delta 번역
# ──────────────────────────────────────────────────────────
with tab_delta:
    st.header("🔄 Delta 번역 — 수정분만 재번역")
    st.caption(
        "이전 한글 원문(v1)과 수정된 한글 원문(v2)을 비교해 **달라진 슬라이드만 감지**하여 재번역합니다. "
        "슬라이드가 추가·삭제되어 페이지 번호가 밀려도 제목과 내용을 기반으로 같은 슬라이드를 찾아 매칭합니다. "
        "수정된 슬라이드에는 **UPDATED 배지**, 신규 슬라이드에는 **NEW 배지**가 표시되고, "
        "맨 앞에 변경 전/후 텍스트 비교가 담긴 **변경사항 요약 시트**가 자동으로 추가됩니다."
    )
    st.warning(
        "🚧 **Under Construction** — 현재 테스트 중인 기능입니다. "
        "매칭 정확도는 슬라이드 내 텍스트 양에 따라 달라질 수 있으며, "
        "결과물은 반드시 사람이 검토 후 사용해 주세요. "
        "오작동이나 개선 의견은 **도지영**에게 알려주세요!",
        icon="⚠️"
    )

    dc1, dc2 = st.columns(2)
    with dc1:
        delta_lang    = st.selectbox("번역 언어", ["English", "Japanese", "Chinese"], key="delta_lang")
        delta_min_pt  = st.number_input("최소 허용 폰트 크기 (pt)", min_value=1, max_value=40, value=7, key="delta_min_pt")
    with dc2:
        st.markdown("##### 파일 업로드")
        old_translated_file = st.file_uploader("① 이전 한글 원문 v1 (.pptx)", type=["pptx"], key="old_tr")
        new_original_file   = st.file_uploader("② 수정된 한글 원문 v2 (.pptx)", type=["pptx"], key="new_orig")

    delta_api_key = st.secrets.get("ANTHROPIC_API_KEY", "")
    if not delta_api_key:
        st.error("⚠️ API Key 미설정.")

    if old_translated_file and new_original_file and delta_api_key:
        st.success(f"✅ **{old_translated_file.name}** (v1) + **{new_original_file.name}** (v2) 업로드 완료")

        if st.button("🔍 Delta 감지 후 번역", type="primary", use_container_width=True):
            import time as _time

            start_time = _time.time()
            delta_lang_str  = LANG_MAP.get(delta_lang, "English")
            active_glossary = get_active_glossary()
            client          = anthropic.Anthropic(api_key=delta_api_key)

            prs_old = Presentation(io.BytesIO(old_translated_file.read()))
            prs_new = Presentation(io.BytesIO(new_original_file.read()))

            # ── Delta 감지 (콘텐츠 지문 기반 매칭) ───────────
            with st.spinner("🔍 변경된 슬라이드 감지 중..."):

                def _get_slide_tokens(slide):
                    """슬라이드 텍스트 줄 리스트 — text_frame + table 모두 포함"""
                    lines = []
                    for shape in slide.shapes:
                        # 일반 텍스트 박스
                        if getattr(shape, "has_text_frame", False) and shape.text_frame:
                            for para in shape.text_frame.paragraphs:
                                t = para.text.strip()
                                if t:
                                    lines.append(t)
                        # 표(Table) 안 텍스트
                        if getattr(shape, "has_table", False) and shape.table:
                            for row in shape.table.rows:
                                for cell in row.cells:
                                    if getattr(cell, "text_frame", None):
                                        for para in cell.text_frame.paragraphs:
                                            t = para.text.strip()
                                            if t:
                                                lines.append(t)
                    return lines

                def _get_title(lines):
                    """첫 번째 비어있지 않은 줄을 제목으로 사용"""
                    return lines[0] if lines else ""

                def _similarity(lines_a, lines_b):
                    """제목 가중 Jaccard 유사도
                    - 빈 슬라이드 둘 다: 0.0 반환 (빈 슬라이드끼리 매칭 방지)
                    - 제목이 같으면 +0.3 보너스 (제목 가중치)
                    - 나머지는 단어 단위 Jaccard
                    """
                    if not lines_a or not lines_b:
                        return 0.0   # 빈 슬라이드는 매칭 안 함

                    set_a = set(" ".join(lines_a).split())
                    set_b = set(" ".join(lines_b).split())
                    inter = len(set_a & set_b)
                    union = len(set_a | set_b)
                    jaccard = inter / union if union else 0.0

                    # 제목 일치 보너스
                    title_a = _get_title(lines_a)
                    title_b = _get_title(lines_b)
                    title_bonus = 0.3 if title_a and title_b and title_a == title_b else 0.0

                    return min(jaccard + title_bonus, 1.0)

                def _text_changed(lines_a, lines_b):
                    """실제 텍스트 내용이 바뀌었는지 판단
                    단어 집합이 완전히 동일하면 unchanged, 아니면 changed
                    → 숫자/날짜 한 글자만 바뀌어도 감지"""
                    return " ".join(sorted(lines_a)) != " ".join(sorted(lines_b))

                # v1, v2 슬라이드 텍스트 추출
                v1_slides = [_get_slide_tokens(s) for s in prs_old.slides]
                v2_slides = [_get_slide_tokens(s) for s in prs_new.slides]

                MATCH_THRESHOLD = 0.35   # 이 이상이면 "같은 슬라이드" 후보

                # ── Hungarian-style 최적 매칭 ──────────────────
                # 모든 (v2, v1) 쌍의 유사도를 먼저 계산한 뒤
                # 높은 점수 순으로 탐욕적으로 할당 → Greedy 순서 의존성 해결
                score_pairs = []
                for v2_idx, v2_lines in enumerate(v2_slides):
                    for v1_idx, v1_lines in enumerate(v1_slides):
                        score = _similarity(v1_lines, v2_lines)
                        if score >= MATCH_THRESHOLD:
                            score_pairs.append((score, v2_idx, v1_idx))

                # 유사도 높은 순 정렬 후 1:1 매칭
                score_pairs.sort(reverse=True)
                used_v2 = set()
                used_v1 = set()
                match_map = {v2_idx: None for v2_idx in range(len(v2_slides))}

                for score, v2_idx, v1_idx in score_pairs:
                    if v2_idx in used_v2 or v1_idx in used_v1:
                        continue
                    match_map[v2_idx] = v1_idx
                    used_v2.add(v2_idx)
                    used_v1.add(v1_idx)

                # 변경/신규 슬라이드 분류
                total_new       = len(v2_slides)
                changed_indices = []
                new_indices     = []
                delta_report    = {}

                for v2_idx in range(total_new):
                    v1_idx   = match_map[v2_idx]
                    v2_lines = v2_slides[v2_idx]

                    if v1_idx is None:
                        # 신규 슬라이드
                        new_indices.append(v2_idx)
                        delta_report[v2_idx] = {
                            "before": [], "after": v2_lines,
                            "v1_idx": None, "type": "new"
                        }
                    else:
                        v1_lines = v1_slides[v1_idx]
                        if _text_changed(v1_lines, v2_lines):
                            # 매칭됐는데 내용 다름 → 수정
                            changed_indices.append(v2_idx)
                            delta_report[v2_idx] = {
                                "before": v1_lines, "after": v2_lines,
                                "v1_idx": v1_idx, "type": "modified"
                            }
                        # 완전 동일 → 무시

            all_delta_indices = changed_indices + new_indices
            if not all_delta_indices:
                st.info("✅ 변경된 슬라이드가 없습니다. 번역이 필요하지 않아요!")
            else:
                summary_parts = []
                if changed_indices:
                    summary_parts.append(f"수정 **{len(changed_indices)}개** → Slide {[i+1 for i in changed_indices]}")
                if new_indices:
                    summary_parts.append(f"신규 **{len(new_indices)}개** → Slide {[i+1 for i in new_indices]}")
                st.info(f"🔄 변경 감지: {' · '.join(summary_parts)}")

                # ── 앱 화면: 변경 전/후 텍스트 리포트 ────────
                with st.expander("📋 변경사항 미리보기 (변경 전/후)", expanded=True):
                    for si in all_delta_indices:
                        info   = delta_report[si]
                        v1_ref = f" ← v1 Slide {info['v1_idx']+1}" if info["v1_idx"] is not None else ""
                        badge  = "🆕 신규" if info["type"] == "new" else "🔸 수정"
                        st.markdown(f"**{badge} Slide {si+1}**{v1_ref}")
                        before_lines = info["before"]
                        after_lines  = info["after"]

                        # 내용 기반 diff — 추가/삭제/유지를 텍스트 집합으로 판단
                        before_set = set(before_lines)
                        after_set  = set(after_lines)
                        added      = after_set  - before_set   # v2에만 있는 줄
                        removed    = before_set - after_set    # v1에만 있는 줄
                        kept       = before_set & after_set    # 동일한 줄

                        rows = []
                        # 유지된 줄
                        for t in before_lines:
                            if t in kept:
                                rows.append({"변경": "", "변경 전 (v1)": t, "변경 후 (v2)": t})
                        # 삭제된 줄 (v1에만)
                        for t in before_lines:
                            if t in removed:
                                rows.append({"변경": "🔴 삭제", "변경 전 (v1)": t, "변경 후 (v2)": ""})
                        # 추가된 줄 (v2에만)
                        for t in after_lines:
                            if t in added:
                                rows.append({"변경": "🟢 추가", "변경 전 (v1)": "", "변경 후 (v2)": t})

                        if not rows:
                            rows = [{"변경": "🆕 신규 슬라이드", "변경 전 (v1)": "", "변경 후 (v2)": " / ".join(after_lines[:5])}]

                        st.dataframe(rows, use_container_width=True, hide_index=True)
                        st.write("")

                progress_bar = st.progress(0)
                status_text  = st.empty()
                log_area     = st.empty()
                log_lines    = []

                for step, slide_idx in enumerate(all_delta_indices):
                    slide  = prs_new.slides[slide_idx]
                    texts  = get_slide_texts(slide)
                    progress_bar.progress((step + 1) / len(all_delta_indices))

                    if not texts:
                        log_lines.append(f"⏭️  Slide {slide_idx+1:2d} — 한국어 없음")
                        log_area.code("\n".join(log_lines))
                        continue

                    slide_type = detect_slide_type(texts)
                    status_text.text(f"Slide {slide_idx+1}/{total_new} [{slide_type}] 번역 중...")

                    for attempt in range(2):
                        try:
                            translated_map = translate_slide(
                                client, texts, slide_type, delta_lang_str, active_glossary
                            )
                            break
                        except Exception as e:
                            if attempt == 0:
                                log_lines.append(f"⚠️  Slide {slide_idx+1:2d} 재시도... ({e})")
                                log_area.code("\n".join(log_lines))
                                _time.sleep(2)
                            else:
                                translated_map = {}
                                log_lines.append(f"❌  Slide {slide_idx+1:2d} 실패 — 원본 유지")

                    # 텍스트 교체
                    shape_map = {}
                    def _collect(shapes):
                        for s in shapes:
                            shape_map[s.shape_id] = s
                            if getattr(s, "shape_type", None) == 6:
                                _collect(s.shapes)
                    _collect(slide.shapes)

                    for ti, text_info in enumerate(texts):
                        tr = translated_map.get(str(ti))
                        if isinstance(tr, dict):
                            tr = tr.get("text") or tr.get("translation") or ""
                        if not tr or not isinstance(tr, str):
                            continue
                        tr = tr.strip()
                        if not tr:
                            continue
                        shape = shape_map.get(text_info["shape_id"])
                        if shape is None:
                            continue
                        para = None
                        if getattr(shape, "has_text_frame", False) and shape.text_frame:
                            paras = shape.text_frame.paragraphs
                            if text_info["para_idx"] < len(paras):
                                para = paras[text_info["para_idx"]]
                        elif getattr(shape, "has_table", False) and shape.table:
                            all_paras = []
                            for row in shape.table.rows:
                                for cell in row.cells:
                                    if getattr(cell, "text_frame", None):
                                        all_paras.extend(cell.text_frame.paragraphs)
                            if text_info["para_idx"] < len(all_paras):
                                para = all_paras[text_info["para_idx"]]
                        if para is not None:
                            replace_para_text(para, tr, shape=shape, min_pt=delta_min_pt)

                    # 배지 추가 — 수정/신규 구분
                    from pptx.util import Inches, Pt as _Pt
                    from pptx.dml.color import RGBColor as _RGB
                    from pptx.enum.text import PP_ALIGN as _ALIGN
                    is_new      = delta_report.get(slide_idx, {}).get("type") == "new"
                    badge_text  = "NEW" if is_new else "UPDATED"
                    badge_color = _RGB(0x2E, 0x7D, 0x32) if is_new else _RGB(0xFF, 0xC0, 0x00)
                    border_color= _RGB(0x1B, 0x5E, 0x20) if is_new else _RGB(0xCC, 0x99, 0x00)
                    text_color  = _RGB(0xFF, 0xFF, 0xFF) if is_new else _RGB(0x4A, 0x2D, 0x00)

                    badge_w, badge_h = Inches(1.0), Inches(0.30)
                    badge_l = prs_new.slide_width - badge_w - Inches(0.18)
                    badge_t = Inches(0.12)
                    badge_shape = slide.shapes.add_shape(1, badge_l, badge_t, badge_w, badge_h)
                    badge_shape.fill.solid()
                    badge_shape.fill.fore_color.rgb = badge_color
                    badge_shape.line.color.rgb      = border_color
                    badge_shape.line.width          = _Pt(0.75)
                    tf = badge_shape.text_frame
                    tf.word_wrap = False
                    p  = tf.paragraphs[0]
                    p.alignment = _ALIGN.CENTER
                    run = p.add_run()
                    run.text           = badge_text
                    run.font.bold      = True
                    run.font.size      = _Pt(8.5)
                    run.font.color.rgb = text_color

                    log_lines.append(f"✅ Slide {slide_idx+1:2d} [{slide_type:15s}] — {len(texts)}개 번역")
                    log_area.code("\n".join(log_lines))
                    _time.sleep(0.3)

                # ── 변경사항 요약 시트 (맨 앞 삽입) ──────────
                status_text.text("📋 변경사항 요약 시트 생성 중...")
                blank_layout  = prs_new.slide_layouts[6]
                summary_slide = prs_new.slides.add_slide(blank_layout)
                xml_slides    = prs_new.slides._sldIdLst
                new_entry     = xml_slides[-1]
                xml_slides.remove(new_entry)
                xml_slides.insert(0, new_entry)

                sw = prs_new.slide_width

                def _tb(slide, text, l, t, w, h, fs=13, bold=False,
                        color=_RGB(30,30,30), align=_ALIGN.LEFT):
                    tb = slide.shapes.add_textbox(l, t, w, h)
                    tf = tb.text_frame; tf.word_wrap = True
                    p  = tf.paragraphs[0]; p.alignment = align
                    r  = p.add_run(); r.text = text
                    r.font.size = _Pt(fs); r.font.bold = bold
                    r.font.color.rgb = color

                hdr = summary_slide.shapes.add_shape(1, 0, 0, sw, Inches(1.05))
                hdr.fill.solid(); hdr.fill.fore_color.rgb = _RGB(0x1B, 0x3A, 0x6B)
                hdr.line.fill.background()
                _tb(summary_slide, "Translation Update — Change Summary",
                    Inches(0.4), Inches(0.22), sw - Inches(0.8), Inches(0.65),
                    fs=20, bold=True, color=_RGB(255,255,255))

                stats  = [("Total Slides", total_new), ("수정", len(changed_indices)), ("신규", len(new_indices))]
                c_fill = [_RGB(0x1B,0x3A,0x6B), _RGB(0xFF,0xC0,0x00), _RGB(0x2E,0x7D,0x32)]
                c_text = [_RGB(255,255,255), _RGB(0x33,0x1A,0x00), _RGB(255,255,255)]
                cw, ch = Inches(2.0), Inches(1.0)
                for i, (lbl, val) in enumerate(stats):
                    cl = Inches(0.4) + i * (cw + Inches(0.25))
                    ct = Inches(1.35)
                    card = summary_slide.shapes.add_shape(1, cl, ct, cw, ch)
                    card.fill.solid(); card.fill.fore_color.rgb = c_fill[i]
                    card.line.fill.background()
                    _tb(summary_slide, str(val),
                        cl+Inches(0.1), ct+Inches(0.04), cw-Inches(0.2), Inches(0.55),
                        fs=28, bold=True, color=c_text[i], align=_ALIGN.CENTER)
                    _tb(summary_slide, lbl,
                        cl+Inches(0.1), ct+Inches(0.58), cw-Inches(0.2), Inches(0.35),
                        fs=10, color=c_text[i], align=_ALIGN.CENTER)

                nums_str = ", ".join([f"Slide {i+2}" for i in all_delta_indices])
                _tb(summary_slide, f"변경/신규 slides:  {nums_str}",
                    Inches(0.4), Inches(2.55), sw-Inches(0.8), Inches(0.4), fs=11)
                _tb(summary_slide,
                    "Slides marked with UPDATED badge (top-right) contain revised translations.",
                    Inches(0.4), Inches(3.05), sw-Inches(0.8), Inches(0.4),
                    fs=11, color=_RGB(0x55,0x55,0x55))

                # ── 슬라이드별 변경 전/후 텍스트 블록 ────────
                y_cursor = Inches(3.6)
                for si in all_delta_indices:
                    info         = delta_report[si]
                    before_lines = info["before"]
                    after_lines  = info["after"]
                    label        = "🆕 NEW" if info["type"] == "new" else "🔸 UPDATED"
                    v1_ref       = f" (← v1 Slide {info['v1_idx']+1})" if info["v1_idx"] is not None else ""

                    _tb(summary_slide, f"{label}  Slide {si+2}{v1_ref}",
                        Inches(0.4), y_cursor, Inches(4.0), Inches(0.28),
                        fs=10, bold=True, color=_RGB(0x1B,0x3A,0x6B))
                    y_cursor += Inches(0.3)

                    before_text = " / ".join(before_lines[:6]) if before_lines else "(없음)"
                    _tb(summary_slide, f"[전] {before_text}",
                        Inches(0.5), y_cursor, sw - Inches(1.0), Inches(0.32),
                        fs=9, color=_RGB(0x99,0x99,0x99))
                    y_cursor += Inches(0.33)

                    after_text = " / ".join(after_lines[:6]) if after_lines else "(없음)"
                    _tb(summary_slide, f"[후] {after_text}",
                        Inches(0.5), y_cursor, sw - Inches(1.0), Inches(0.32),
                        fs=9, bold=True, color=_RGB(0x1B,0x3A,0x6B))
                    y_cursor += Inches(0.42)

                    if y_cursor > Inches(6.8):
                        remaining = len(all_delta_indices) - all_delta_indices.index(si) - 1
                        _tb(summary_slide, f"... 외 {remaining}개 슬라이드",
                            Inches(0.5), y_cursor, sw - Inches(1.0), Inches(0.3),
                            fs=9, color=_RGB(0x99,0x99,0x99))
                        break

                # ── 저장 & 다운로드 ────────────────────────
                output = io.BytesIO()
                prs_new.save(output)
                output.seek(0)
                elapsed  = _time.time() - start_time
                out_name = new_original_file.name.rsplit(".", 1)[0] + "_EN_delta.pptx"

                progress_bar.progress(1.0)
                status_text.empty()
                st.success(f"🎉 Delta 번역 완료! — 수정 {len(changed_indices)}개 · 신규 {len(new_indices)}개 재번역 (⏱️ {elapsed:.1f}초)")
                st.download_button(
                    label=f"⬇️ {out_name} 다운로드",
                    data=output,
                    file_name=out_name,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True,
                    type="primary",
                )


with tab_glossary:
    st.header("📖 Glossary")
    st.info("💬 Glossary 추가/수정 요청은 **도지영** (michelle@krafton.com)에게 연락해주세요. [Slack](https://krafton.enterprise.slack.com/team/U02RWGEGJ5B)")

    active   = get_active_glossary()

    search = st.text_input("🔍 검색", placeholder="한국어 또는 영문으로 검색...")
    
    db_data = load_glossary_db()
    approved = db_data.get('approved_glossary', {})
    st.caption(f"총 **{len(active)}개** 항목 (기본 {len(BASE_GLOSSARY)} + 추가 {len(approved)}개)")

    NAME_KEYS = set(list(BASE_GLOSSARY.keys())[:39])

    col_a, col_b = st.columns(2)
    with col_a:
        st.markdown("#### 👤 인명")
        for ko, en in BASE_GLOSSARY.items():
            if ko not in NAME_KEYS:
                continue
            if search and search.lower() not in ko.lower() and search.lower() not in en.lower():
                continue
            st.markdown(f"`{ko}` → {en}")

    with col_b:
        st.markdown("#### 📝 용어")
        for ko, en in BASE_GLOSSARY.items():
            if ko in NAME_KEYS:
                continue
            if search and search.lower() not in ko.lower() and search.lower() not in en.lower():
                continue
            st.markdown(f"`{ko}` → {en}")

    if approved:
        st.divider()
        st.markdown(f"#### ✅ 추가 등록된 항목 — {len(approved)}개")
        for ko, en in approved.items():
            if search and search.lower() not in ko.lower() and search.lower() not in en.lower():
                continue
            st.markdown(f"`{ko}` → {en}")
